from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG      = RGBColor(0x0f, 0x11, 0x17)
SURFACE = RGBColor(0x1a, 0x1d, 0x27)
BORDER  = RGBColor(0x2a, 0x2d, 0x3a)
ACCENT  = RGBColor(0x6c, 0x8e, 0xf7)
TEXT    = RGBColor(0xe2, 0xe4, 0xf0)
MUTED   = RGBColor(0x7a, 0x7d, 0x90)
GREEN   = RGBColor(0x4a, 0xde, 0x80)
YELLOW  = RGBColor(0xfb, 0xbf, 0x24)
DARK2   = RGBColor(0x1e, 0x22, 0x35)
DARK3   = RGBColor(0x10, 0x18, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_slide():
    sl = prs.slides.add_slide(blank)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl


def box(sl, x, y, w, h, fill=SURFACE, border=BORDER, bw=Pt(1)):
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = bw
    return shp


def txt(sl, text, x, y, w, h, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def pill(sl, text, x, y, w=1.8, h=0.3, fc=DARK2, tc=ACCENT, ts=8):
    shp = box(sl, x, y, w, h, fill=fc, border=BORDER, bw=Pt(0.5))
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(ts)
    run.font.color.rgb = tc
    return shp


def node(sl, text, x, y, w, h, fc=DARK2, tc=ACCENT, ts=9):
    shp = box(sl, x, y, w, h, fill=fc, border=ACCENT, bw=Pt(0.8))
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(ts)
    run.font.bold = True
    run.font.color.rgb = tc
    return shp


def header_bar(sl, tag, sub=""):
    box(sl, 0, 0, 13.33, 0.72, fill=SURFACE, border=BORDER)
    txt(sl, "LMS AI Copilot", 0.3, 0.16, 2.6, 0.42, size=11, bold=True, color=TEXT)
    pill(sl, tag, 3.1, 0.2, 2.2, 0.32, fc=DARK2, tc=ACCENT, ts=9)
    if sub:
        txt(sl, sub, 5.5, 0.22, 6.0, 0.32, size=8.5, color=MUTED)
    txt(sl, "ISD Team  2026", 11.7, 0.22, 1.5, 0.32, size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)


def table_box(sl, name, fields, x, y, w, h):
    hdr = box(sl, x, y, w, 0.3, fill=DARK2, border=ACCENT, bw=Pt(1))
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = name
    r.font.size = Pt(8.5)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    bdy = box(sl, x, y + 0.3, w, h - 0.3, fill=SURFACE, border=BORDER)
    tf2 = bdy.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = fields
    r2.font.size = Pt(7)
    r2.font.color.rgb = MUTED


def layer_box(sl, name, items, x, y, w, h, fc=DARK2):
    hdr = box(sl, x, y, w, 0.32, fill=DARK2, border=ACCENT, bw=Pt(1))
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = name
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = ACCENT

    bdy = box(sl, x, y + 0.32, w, h - 0.32, fill=fc, border=BORDER)
    tf2 = bdy.text_frame
    tf2.word_wrap = True
    for j, item in enumerate(items):
        p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "· " + item
        r2.font.size = Pt(8)
        r2.font.color.rgb = MUTED


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
sl = add_slide()
box(sl, 1.5, 3.52, 10.33, 0.05, fill=ACCENT, border=ACCENT)

txt(sl, "LMS-Aware AI Copilot", 0, 2.2, 13.33, 1.0,
    size=44, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
txt(sl, "for Personalized Learning", 0, 3.1, 13.33, 0.65,
    size=28, color=ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "HUFS  ·  ISD Team  ·  June 2026", 0, 3.75, 13.33, 0.45,
    size=13, color=MUTED, align=PP_ALIGN.CENTER)

pills = ["Chrome Extension MV3", "FastAPI", "SQLite FTS5",
         "Groq llama-3.3-70b", "RAG Pipeline", "SSE Streaming"]
for i, p_text in enumerate(pills):
    pill(sl, p_text, 0.55 + i * 2.05, 4.5, 1.92, 0.34, tc=ACCENT, ts=9)

txt(sl, "5-min presentation  ·  Live Demo  ·  Q&A", 0, 5.25, 13.33, 0.4,
    size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ERD
# ══════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "ERD", "SQLite · 9 tables · FTS5 virtual tables")
box(sl, 0.2, 0.8, 12.93, 6.55, fill=SURFACE, border=BORDER)
txt(sl, "Data Model", 0.5, 0.88, 4, 0.32, size=11, bold=True, color=ACCENT)

tables = [
    ("User",              "user_id PK\nlms_id\nenc_session_cookie\nscraping_interval\nlast_sync_at",    0.35, 1.28, 2.05, 1.38),
    ("Course",            "course_id PK\nlms_url_id UK\ncourse_code\ntitle",                            0.35, 2.85, 2.05, 1.05),
    ("Enrollment",        "enroll_id PK\nuser_id FK\ncourse_id FK\nrole",                               0.35, 4.1,  2.05, 1.0),
    ("Material",          "material_id PK\ncourse_id FK\ntitle · file_type\nchecksum",                  2.6,  1.28, 2.1,  1.05),
    ("Doc_Chunk",         "chunk_id PK\nmaterial_id FK\ncontent · page_ref\nchunk_index",               2.6,  2.55, 2.1,  1.05),
    ("Learning_Activity", "activity_id PK\nuser_id FK\ncourse_id FK\ntitle · status · due_date",        2.6,  3.82, 2.1,  1.05),
    ("Personal_Log",      "log_id PK\nuser_id FK\ncourse_id FK\nmaterial_id FK\nstay_time",             5.0,  1.28, 2.1,  1.25),
    ("Chat_Session",      "session_id PK\nuser_id FK\ncourse_id FK\ncreated_at",                        5.0,  2.75, 2.1,  1.0),
    ("Chat_Log",          "chat_id PK\nsession_id FK\nrole · content\nkeywords · sources\nfeedback_score", 5.0, 3.97, 2.1, 1.35),
]
for name, fields, x, y, w, h in tables:
    table_box(sl, name, fields, x, y, w, h)

pill(sl, "Doc_Chunk_fts  (FTS5 index)", 7.4, 2.55, 2.4, 0.3, fc=RGBColor(0x0f, 0x20, 0x15), tc=GREEN, ts=8)
pill(sl, "Chat_Log_fts  (FTS5 index)",  7.4, 3.1,  2.4, 0.3, fc=RGBColor(0x0f, 0x20, 0x15), tc=GREEN, ts=8)

txt(sl, "Key Relationships", 7.35, 1.28, 5.4, 0.3, size=9, bold=True, color=ACCENT)
rels = [
    "User ──< Enrollment >── Course",
    "Course ──< Material ──< Doc_Chunk",
    "User  ──< Learning_Activity",
    "User  ──< Personal_Log",
    "User  ──< Chat_Session ──< Chat_Log",
    "Material ──< Personal_Log (viewed)",
]
for i, r in enumerate(rels):
    txt(sl, r, 7.35, 1.62 + i * 0.33, 5.5, 0.3, size=8, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "System Architecture", "6-layer end-to-end")
box(sl, 0.2, 0.8, 12.93, 6.55, fill=SURFACE, border=BORDER)
txt(sl, "End-to-End System", 0.5, 0.88, 5, 0.32, size=11, bold=True, color=ACCENT)

top_layers = [
    ("Chrome Extension",  ["Sidebar UI (panel.html)", "Cookie Interceptor", "JSESSIONID auto-capture", "Personal_Log stay-time"],     0.35, 1.28, 2.3, 2.1),
    ("FastAPI Backend",   ["POST /sync · /sync/delta", "POST /chat · /chat/stream(SSE)", "GET /chat/{id} · /courses", "POST /log · /feedback"], 2.85, 1.28, 2.5, 2.1),
    ("e-Class Scraper",   ["main_form.acl (auth)", "efile_download2 (PDF/PPTX)", "report_list.acl POST", "project_list.acl POST", "test_list_form GET"], 5.55, 1.28, 2.3, 2.1),
    ("Ingest Pipeline",   ["chunker.py", "PDF · PPTX → text extract", "Page-level chunking", "MD5 checksum dedup", "FTS5 index insert"], 8.05, 1.28, 2.3, 2.1),
]
for name, items, x, y, w, h in top_layers:
    layer_box(sl, name, items, x, y, w, h, fc=DARK3)

bot_layers = [
    ("SQLite Database", ["User · Course · Enrollment", "Material · Doc_Chunk", "Doc_Chunk_fts (FTS5)", "Learning_Activity", "Chat_Session · Chat_Log + fts", "Personal_Log"], 0.35, 3.65, 3.5, 2.55),
    ("RAG Engine",      ["Smart Router — ACTIVITY vs MATERIAL keywords", "Keyword Extractor — llama-3.1-8b (token saving)", "FTS5 Doc_Chunk MATCH LIMIT 8", "Hybrid History — recent 4 + FTS relevant 3", "Prompt Builder + D-day tag (Python)", "Groq llama-3.3-70b → 8b fallback on 429", "SSE StreamingResponse — real-time render"], 4.05, 3.65, 6.3, 2.55),
    ("Notification Sync", ["POST /sync/delta", "e-Class alert triggered", "Changed courses only", "Incremental update"], 10.55, 3.65, 2.3, 2.55),
]
for name, items, x, y, w, h in bot_layers:
    layer_box(sl, name, items, x, y, w, h, fc=DARK3)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SEQUENCE DIAGRAMS
# ══════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Sequence Diagrams", "LMS Sync · RAG Chat")
box(sl, 0.2, 0.8, 12.93, 6.55, fill=SURFACE, border=BORDER)

# divider
box(sl, 6.6, 0.85, 0.04, 6.4, fill=BORDER, border=BORDER)

# ── LEFT: Sync ──
txt(sl, "① LMS Sync", 0.4, 0.9, 4, 0.3, size=10, bold=True, color=ACCENT)

sync_actors = [("User", 0.45), ("Ext", 1.55), ("API", 2.65), ("Scraper", 3.7), ("e-Class", 4.75), ("SQLite", 5.75)]
for name, x in sync_actors:
    node(sl, name, x, 1.28, 0.85, 0.3, ts=8)
    box(sl, x + 0.4, 1.58, 0.05, 5.0, fill=BORDER, border=BORDER)

sync_steps = [
    (0.45, 1.55, "동기화 클릭",            1.85),
    (1.55, 2.65, "POST /sync {cookie}",    2.15),
    (2.65, 3.7,  "run_sync()",             2.45),
    (3.7,  4.75, "GET main_form.acl",      2.75),
    (4.75, 5.75, "upsert user/course",     3.05),
    (3.7,  4.75, "POST report_list.acl",   3.35),
    (3.7,  4.75, "POST project_list.acl",  3.6),
    (3.7,  5.75, "clear + upsert activity",3.85),
    (3.7,  5.75, "chunk PDF → FTS5",       4.1),
    (2.65, 5.75, "touch_sync",             4.38),
    (1.55, 0.45, "방금 동기화",            4.65),
]
for x1, x2, label, y in sync_steps:
    cx1 = x1 + 0.425
    cx2 = x2 + 0.425
    arrow = sl.shapes.add_connector(1, Inches(cx1), Inches(y), Inches(cx2), Inches(y))
    arrow.line.color.rgb = ACCENT
    arrow.line.width = Pt(0.8)
    mx = min(cx1, cx2) + 0.05
    txt(sl, label, mx, y - 0.17, abs(cx2 - cx1) + 0.1, 0.18, size=6, color=MUTED)

# ── RIGHT: RAG Chat ──
txt(sl, "② RAG Chat  (SSE Streaming)", 6.85, 0.9, 5.5, 0.3, size=10, bold=True, color=ACCENT)

rag_actors = [("User", 6.9), ("Ext", 7.95), ("API", 9.05), ("SQLite", 10.1), ("Groq", 11.15)]
for name, x in rag_actors:
    node(sl, name, x, 1.28, 0.85, 0.3, ts=8)
    box(sl, x + 0.4, 1.58, 0.05, 5.0, fill=BORDER, border=BORDER)

rag_steps = [
    (6.9,  7.95, "질문 입력",                  1.85),
    (7.95, 9.05, "POST /chat/stream",           2.15),
    (9.05, 10.1, "Smart Router",                2.45),
    (9.05, 10.1, "get_hybrid_history",          2.72),
    (9.05, 11.15,"extract_keywords (8b)",       2.99),
    (9.05, 10.1, "FTS5 MATCH LIMIT 8",          3.26),
    (9.05, 11.15,"llama-3.3-70b stream=True",   3.53),
    (11.15, 7.95,"← SSE delta chunks",          3.95),
    (7.95, 6.9,  "← 실시간 렌더링",             4.28),
    (9.05, 10.1, "INSERT Chat_Log + FTS",        4.58),
    (9.05, 7.95, "done {session_id, sources}",   4.85),
    (6.9,  9.05, "POST /feedback ±1",           5.2),
]
for x1, x2, label, y in rag_steps:
    cx1 = x1 + 0.425
    cx2 = x2 + 0.425
    arrow = sl.shapes.add_connector(1, Inches(cx1), Inches(y), Inches(cx2), Inches(y))
    arrow.line.color.rgb = ACCENT
    arrow.line.width = Pt(0.8)
    mx = min(cx1, cx2) + 0.05
    txt(sl, label, mx, y - 0.17, abs(cx2 - cx1) + 0.1, 0.18, size=6, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DEMO
# ══════════════════════════════════════════════════════════════════════════
sl = add_slide()
header_bar(sl, "Live Demo", "Chrome Extension  ·  FastAPI  ·  Groq")
box(sl, 0.2, 0.8, 12.93, 6.55, fill=SURFACE, border=BORDER)
txt(sl, "Demo Scenario", 0.5, 0.88, 4, 0.32, size=11, bold=True, color=ACCENT)

steps_demo = [
    ("① Sync",           "학번 입력 → 동기화 → e-Class 전 과목 자동 스크래핑 (PDF 청크화 + 과제 수집)"),
    ("② Course Detect",  "강의실 클릭 → URL KJKEY 자동 감지 → 헤더 과목명 자동 표시"),
    ("③ Assignment Q&A", '"이번 주 과제 마감일이 언제야?" → D-day 태그 포함 스트리밍 답변'),
    ("④ Material Q&A",   '"Merge Sort 시간복잡도 설명해줘" → FTS5 검색 → 출처 페이지 태그 표시'),
    ("⑤ Session Reuse",  '"그럼 Quick Sort랑 비교하면?" → 세션 재사용, 이전 맥락 유지'),
    ("⑥ Feedback",       "👍 도움됐어요 클릭 → DB Chat_Log.feedback_score = 1 저장"),
]
for i, (title, desc) in enumerate(steps_demo):
    y = 1.35 + i * 0.82
    node(sl, title, 0.45, y, 2.0, 0.35, tc=ACCENT, ts=9)
    txt(sl, desc, 2.65, y + 0.02, 10.2, 0.35, size=9, color=TEXT)

tech_items = [
    ("Smart Routing",      ACCENT),
    ("D-day Python calc",  ACCENT),
    ("SSE Streaming",      GREEN),
    ("Hybrid FTS History", ACCENT),
    ("Source Page Tags",   YELLOW),
    ("Personal_Log",       MUTED),
]
for i, (t, c) in enumerate(tech_items):
    pill(sl, t, 0.45 + i * 2.08, 6.55, 1.98, 0.3, tc=c, ts=8)


prs.save("ISD_Team_Presentation.pptx")
print("saved -> ISD_Team_Presentation.pptx")
